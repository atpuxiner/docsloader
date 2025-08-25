import logging
from typing import AsyncGenerator, Generator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pathlib import Path

from docsloader.base import BaseLoader, DocsData
from docsloader.utils import format_table, format_image

logger = logging.getLogger(__name__)


class PptxLoader(BaseLoader):

    async def load_by_basic(self) -> AsyncGenerator[DocsData, None]:
        for item in self.extract_by_python_pptx(tmpfile=self.tmpfile):
            self.metadata.update(
                page=item.get("page"),
                page_total=item.get("page_total"),
            )
            yield DocsData(
                type=item.get("type"),
                text=item.get("text"),
                data=item.get("data"),
                metadata=self.metadata,
            )

    @staticmethod
    def extract_by_python_pptx(tmpfile: str) -> Generator[dict, None, None]:
        presentation = Presentation(tmpfile)
        image_dir = Path(f"{tmpfile}.images")
        image_dir.mkdir(parents=True, exist_ok=True)
        page_total = len(presentation.slides)
        for slide_idx, slide in enumerate(presentation.slides):
            logger.debug(f"Processing slide {slide_idx + 1}")
            for shape_idx, shape in enumerate(slide.shapes):
                extracted_data = PptxLoader.extract_shape(
                    shape=shape,
                    image_dir=image_dir,
                    image_idx=f"{slide_idx}-{shape_idx}",
                )
                if extracted_data:
                    extracted_data.update(
                        page=slide_idx + 1,
                        page_total=page_total,
                    )
                    yield extracted_data
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub_shape_idx, sub_shape in enumerate(shape.shapes):
                        group_extracted_data = PptxLoader.extract_shape(
                            shape=sub_shape,
                            image_dir=image_dir,
                            image_idx=f"{slide_idx}-{shape_idx}-{sub_shape_idx}",
                        )
                        if group_extracted_data:
                            group_extracted_data.update(
                                page=slide_idx + 1,
                                page_total=page_total,
                            )
                            yield group_extracted_data
        if image_dir.is_dir() and not any(image_dir.iterdir()):
            try:
                image_dir.rmdir()
            except OSError:
                logger.debug(f"Could not remove empty image directory: {image_dir}")

    @staticmethod
    def extract_shape(shape, image_dir: Path, image_idx: str) -> dict:
        """
        解析单个 shape 对象，提取其中的文本、表格和图片信息。
        """
        shape_text = ""
        shape_data = {}
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs).strip()
                if para_text:
                    shape_text += para_text + "\n"
            if shape_text:
                shape_data = {
                    "type": "text",
                    "text": shape_text,
                }
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_title = shape.name if shape.name else "Table"
            table_data = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
            shape_text += f"\n## {table_title}\n"
            shape_text += format_table(table_data)
            shape_data = {
                "type": "table",
                "text": shape_text,
                "data": table_data
            }
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_filename = f"image_{image_idx}.{image.ext}"
            image_path = str(image_dir / image_filename)
            with open(image_path, "wb") as f:
                f.write(image.blob)
            shape_data = {
                "type": "image",
                "text": format_image(image_path),
                "data": image_path,
            }
        return shape_data
